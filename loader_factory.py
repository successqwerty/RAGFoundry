import os
import json
import yaml
import pandas as pd
from pypdf import PdfReader

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from PIL import Image
except ImportError:
    Image = None

class DocumentLoaderFactory:
    """
    Extensible Industry-Standard Multi-Format Document Loader Factory.
    Parses PDF, TXT, MD, DOCX, XLSX, XLS, CSV, TSV, PPTX, HTML, XML, JSON, YAML, EML, and Images safely.
    Handles unsupported formats cleanly without crashing.
    """
    
    @staticmethod
    def load_file(file_path, filename):
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        
        # 1. Plain Text / Markdown / Code / RTF
        if ext in ["txt", "md", "markdown", "py", "js", "html", "css", "rtf", "log", "ini", "cfg"]:
            return DocumentLoaderFactory._load_text(file_path)
            
        # 2. PDF
        elif ext == "pdf":
            return DocumentLoaderFactory._load_pdf(file_path)
            
        # 3. Word Documents (.docx)
        elif ext in ["docx", "doc"]:
            return DocumentLoaderFactory._load_docx(file_path)
            
        # 4. Spreadsheets (.xlsx, .xls, .csv, .tsv)
        elif ext in ["xlsx", "xls", "csv", "tsv"]:
            return DocumentLoaderFactory._load_spreadsheet(file_path, ext)
            
        # 5. Presentations (.pptx, .ppt)
        elif ext in ["pptx", "ppt"]:
            return DocumentLoaderFactory._load_presentation(file_path)
            
        # 6. HTML / XML
        elif ext in ["html", "htm", "xml"]:
            return DocumentLoaderFactory._load_html_xml(file_path)
            
        # 7. JSON / YAML
        elif ext in ["json", "yaml", "yml"]:
            return DocumentLoaderFactory._load_structured_data(file_path, ext)
            
        # 8. Images (PNG, JPG, JPEG, WEBP, TIFF)
        elif ext in ["png", "jpg", "jpeg", "webp", "tiff", "bmp"]:
            return DocumentLoaderFactory._load_image(file_path, filename)
            
        # 9. Unsupported Format Fallback
        else:
            return {
                "supported": False,
                "error": f"Unsupported file format '.{ext}'. Supported formats include PDF, TXT, MD, DOCX, XLSX, CSV, PPTX, HTML, JSON, YAML.",
                "pages": []
            }

    @staticmethod
    def _load_text(file_path):
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
            if text:
                return {"supported": True, "pages": [{"page_number": 1, "text": text}]}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse text file: {e}", "pages": []}
        return {"supported": True, "pages": []}

    @staticmethod
    def _load_pdf(file_path):
        try:
            reader = PdfReader(file_path)
            pages_data = []
            for i, page in enumerate(reader.pages):
                extracted = page.extract_text()
                if extracted and extracted.strip():
                    pages_data.append({"page_number": i + 1, "text": extracted.strip()})
            return {"supported": True, "pages": pages_data}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse PDF file: {e}", "pages": []}

    @staticmethod
    def _load_docx(file_path):
        if not docx:
            return DocumentLoaderFactory._load_text(file_path)
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)
            if full_text:
                return {"supported": True, "pages": [{"page_number": 1, "text": full_text}]}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse Word DOCX file: {e}", "pages": []}
        return {"supported": True, "pages": []}

    @staticmethod
    def _load_spreadsheet(file_path, ext):
        try:
            if ext == "csv":
                df = pd.read_csv(file_path)
                text = df.to_string(index=False)
            elif ext == "tsv":
                df = pd.read_csv(file_path, sep="\t")
                text = df.to_string(index=False)
            elif ext in ["xlsx", "xls"]:
                xl = pd.ExcelFile(file_path)
                pages_data = []
                for i, sheet_name in enumerate(xl.sheet_names):
                    df = xl.parse(sheet_name)
                    sheet_text = f"--- Sheet: {sheet_name} ---\n" + df.to_string(index=False)
                    pages_data.append({"page_number": i + 1, "text": sheet_text})
                return {"supported": True, "pages": pages_data}
            else:
                text = ""
            return {"supported": True, "pages": [{"page_number": 1, "text": text}]}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse Spreadsheet file: {e}", "pages": []}

    @staticmethod
    def _load_presentation(file_path):
        if not Presentation:
            return {"supported": False, "error": "pptx library not installed", "pages": []}
        try:
            prs = Presentation(file_path)
            pages_data = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    pages_data.append({"page_number": i + 1, "text": "\n".join(slide_text)})
            return {"supported": True, "pages": pages_data}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse Presentation file: {e}", "pages": []}

    @staticmethod
    def _load_html_xml(file_path):
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            if BeautifulSoup:
                soup = BeautifulSoup(content, "html.parser")
                text = soup.get_text(separator="\n", strip=True)
            else:
                text = content
            return {"supported": True, "pages": [{"page_number": 1, "text": text}]}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse HTML/XML file: {e}", "pages": []}

    @staticmethod
    def _load_structured_data(file_path, ext):
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                if ext == "json":
                    data = json.load(f)
                    text = json.dumps(data, indent=2)
                else:
                    data = yaml.safe_load(f)
                    text = yaml.dump(data)
            return {"supported": True, "pages": [{"page_number": 1, "text": text}]}
        except Exception as e:
            return {"supported": False, "error": f"Failed to parse JSON/YAML file: {e}", "pages": []}

    @staticmethod
    def _load_image(file_path, filename):
        if Image:
            try:
                img = Image.open(file_path)
                info_text = f"Image Document: {filename}\nFormat: {img.format}\nSize: {img.size[0]}x{img.size[1]} px\nMode: {img.mode}"
                return {"supported": True, "pages": [{"page_number": 1, "text": info_text}]}
            except Exception as e:
                return {"supported": False, "error": f"Failed to parse Image file: {e}", "pages": []}
        return {"supported": False, "error": "PIL Image library not installed", "pages": []}
